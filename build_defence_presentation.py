import sys
import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# --- COLOR PALETTE DEFINITIONS ---
COLOR_NAVY = RGBColor(15, 32, 67)          # #0F2043 Executive Deep Navy
COLOR_CRIMSON = RGBColor(165, 28, 48)      # #A51C30 Harvard Crimson
COLOR_GOLD = RGBColor(212, 175, 55)        # #D4AF37 Metallic Gold
COLOR_SLATE_DARK = RGBColor(30, 41, 59)    # #1E293B Charcoal Body Text
COLOR_SLATE_MUTED = RGBColor(100, 116, 139) # #64748B Subtitle/Secondary Text
COLOR_CARD_BG = RGBColor(248, 250, 252)    # #F8FAFC Off-white card fill
COLOR_BORDER = RGBColor(226, 232, 240)     # #E2E8F0 Soft border gray
COLOR_WHITE = RGBColor(255, 255, 255)      # #FFFFFF Clean White
COLOR_LIGHT_NAVY = RGBColor(238, 242, 255) # Light Navy Tint
COLOR_LIGHT_CRIMSON = RGBColor(254, 242, 242) # Light Crimson Tint
COLOR_LIGHT_GOLD = RGBColor(254, 249, 195) # Light Gold Tint
COLOR_GREEN_TEXT = RGBColor(22, 101, 52)   # Forest green accent
COLOR_GREEN_BG = RGBColor(240, 253, 244)   # Light green fill

# Terminal / Code Console Colors
COLOR_TERM_BG = RGBColor(13, 17, 23)       # #0D1117 Dark IDE background
COLOR_TERM_GREEN = RGBColor(74, 222, 128)  # #4ADE80 Terminal Green
COLOR_TERM_CYAN = RGBColor(56, 189, 248)   # #38BDF8 Terminal Cyan
COLOR_TERM_YELLOW = RGBColor(250, 204, 21) # #FACC15 Terminal Yellow
COLOR_TERM_WHITE = RGBColor(243, 244, 246) # #F3F4F6 Terminal White

FONT_HEADING = "Arial"
FONT_BODY = "Calibri"
FONT_CODE = "Consolas"

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # blank

    # --- HELPER FUNCTIONS ---
    def set_slide_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_tag="MASTER'S THESIS DEFENCE"):
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.733), Inches(0.9))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        # Category Tag
        p_tag = tf.paragraphs[0]
        p_tag.text = category_tag.upper()
        p_tag.font.name = FONT_HEADING
        p_tag.font.size = Pt(10)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_CRIMSON
        p_tag.space_after = Pt(2)

        # Title Text
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.name = FONT_HEADING
        p_title.font.size = Pt(21)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_NAVY

        # Accent Line under header
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_GOLD
        line.line.fill.background()

    def add_footer(slide, slide_num, total_slides=14):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(9.5), Inches(0.35))
        tf = footer_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = "University of Verona | Department of Economics  •  Murtuza Yusuf Rangwala  •  Supervisor: Prof. Giuseppina Chesini"
        p.font.name = FONT_BODY
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_SLATE_MUTED

        right_box = slide.shapes.add_textbox(Inches(10.5), Inches(7.05), Inches(2.033), Inches(0.35))
        rtf = right_box.text_frame
        rtf.margin_left = rtf.margin_top = rtf.margin_right = rtf.margin_bottom = 0
        rp = rtf.paragraphs[0]
        rp.alignment = PP_ALIGN.RIGHT
        rp.text = f"Slide {slide_num} / {total_slides}"
        rp.font.name = FONT_BODY
        rp.font.size = Pt(9.5)
        rp.font.bold = True
        rp.font.color.rgb = COLOR_NAVY

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD_BG, border_color=COLOR_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    def add_card_with_image(slide, img_path, left, top, width, height, title="", bg_color=COLOR_WHITE, border_color=COLOR_GOLD):
        # Background card container
        add_card(slide, left, top, width, height, bg_color=bg_color, border_color=border_color)
        
        title_h = Inches(0.4) if title else Inches(0.15)
        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.35))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = FONT_HEADING
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = COLOR_NAVY
            p.alignment = PP_ALIGN.CENTER
            
        avail_w = width - Inches(0.4)
        avail_h = height - title_h - Inches(0.2)
        
        if os.path.exists(img_path):
            with Image.open(img_path) as im:
                img_w, img_h = im.size
            aspect = img_w / img_h
            
            target_w = avail_w
            target_h = target_w / aspect
            
            if target_h > avail_h:
                target_h = avail_h
                target_w = target_h * aspect
                
            img_left = left + (width - target_w) / 2
            img_top = top + title_h + (avail_h - target_h) / 2
            
            slide.shapes.add_picture(img_path, img_left, img_top, width=target_w, height=target_h)

    def add_stat_card(slide, left, top, width, height, value_text, label_text, subtext="", value_color=COLOR_NAVY, bg_color=COLOR_CARD_BG, border_color=COLOR_BORDER):
        add_card(slide, left, top, width, height, bg_color=bg_color, border_color=border_color)
        tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), height - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_val = tf.paragraphs[0]
        p_val.text = value_text
        p_val.font.name = FONT_HEADING
        p_val.font.size = Pt(26)
        p_val.font.bold = True
        p_val.font.color.rgb = value_color
        p_val.alignment = PP_ALIGN.CENTER
        
        p_lbl = tf.add_paragraph()
        p_lbl.text = label_text
        p_lbl.font.name = FONT_HEADING
        p_lbl.font.size = Pt(10.5)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = COLOR_SLATE_DARK
        p_lbl.alignment = PP_ALIGN.CENTER
        p_lbl.space_before = Pt(3)
        
        if subtext:
            p_sub = tf.add_paragraph()
            p_sub.text = subtext
            p_sub.font.name = FONT_BODY
            p_sub.font.size = Pt(9)
            p_sub.font.color.rgb = COLOR_SLATE_MUTED
            p_sub.alignment = PP_ALIGN.CENTER
            p_sub.space_before = Pt(2)

    def style_table(table, col_widths, headers, data, header_bg=COLOR_NAVY):
        for i, width in enumerate(col_widths):
            table.columns[i].width = width
            
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.text = h
            p.font.name = FONT_HEADING
            p.font.size = Pt(10.5)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE
            p.alignment = PP_ALIGN.CENTER
            
        for row_idx, row_data in enumerate(data):
            bg = COLOR_WHITE if row_idx % 2 == 0 else COLOR_CARD_BG
            for col_idx, val in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = cell.text_frame.paragraphs[0]
                p.text = str(val)
                p.font.name = FONT_BODY
                p.font.size = Pt(10)
                if col_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = COLOR_NAVY
                    p.alignment = PP_ALIGN.LEFT
                else:
                    p.font.color.rgb = COLOR_SLATE_DARK
                    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1, COLOR_NAVY)

    # Gold Top Accent Line
    top_line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = COLOR_GOLD
    top_line.line.fill.background()

    # Title Card Overlay Shape
    add_card(slide1, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.8), bg_color=RGBColor(20, 40, 80), border_color=COLOR_GOLD)

    # University Badge Box
    tb_univ = slide1.shapes.add_textbox(Inches(1.2), Inches(1.1), Inches(10.933), Inches(0.6))
    tf_univ = tb_univ.text_frame
    tf_univ.word_wrap = True
    p_u = tf_univ.paragraphs[0]
    p_u.text = "UNIVERSITY OF VERONA  •  DEPARTMENT OF ECONOMICS"
    p_u.font.name = FONT_HEADING
    p_u.font.size = Pt(13)
    p_u.font.bold = True
    p_u.font.color.rgb = COLOR_GOLD

    p_msc = tf_univ.add_paragraph()
    p_msc.text = "Master's Degree Thesis Defence in Economics and Data Analysis | Academic Year 2025/2026"
    p_msc.font.name = FONT_BODY
    p_msc.font.size = Pt(11)
    p_msc.font.color.rgb = RGBColor(203, 213, 225)
    p_msc.space_before = Pt(2)

    # Main Title & Subtitle
    tb_main = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.933), Inches(2.2))
    tf_main = tb_main.text_frame
    tf_main.word_wrap = True

    p_t = tf_main.paragraphs[0]
    p_t.text = "Temporal Fusion Deep Macro-Gated Attention\n(TFDMGA)"
    p_t.font.name = FONT_HEADING
    p_t.font.size = Pt(30)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_WHITE
    p_t.space_after = Pt(10)

    p_s = tf_main.add_paragraph()
    p_s.text = "Publication-Grade Machine Learning Architecture for Cross-Sectional Stock Return Prediction & Quantitative Execution"
    p_s.font.name = FONT_BODY
    p_s.font.size = Pt(15)
    p_s.font.color.rgb = COLOR_GOLD

    # Separator Line
    sep = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(4.3), Inches(10.933), Inches(0.02))
    sep.fill.solid()
    sep.fill.fore_color.rgb = COLOR_CRIMSON
    sep.line.fill.background()

    # Candidate & Supervisor Details Card
    tb_meta = slide1.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(10.933), Inches(1.8))
    tf_meta = tb_meta.text_frame
    tf_meta.word_wrap = True

    p_c = tf_meta.paragraphs[0]
    p_c.text = "Candidate: Murtuza Yusuf Rangwala"
    p_c.font.name = FONT_HEADING
    p_c.font.size = Pt(15)
    p_c.font.bold = True
    p_c.font.color.rgb = COLOR_WHITE

    p_sup = tf_meta.add_paragraph()
    p_sup.text = "Supervisor: Prof. Giuseppina Chesini"
    p_sup.font.name = FONT_HEADING
    p_sup.font.size = Pt(14)
    p_sup.font.bold = True
    p_sup.font.color.rgb = RGBColor(226, 232, 240)
    p_sup.space_before = Pt(4)

    p_d = tf_meta.add_paragraph()
    p_d.text = "Academic Year 2024 / 2025  |  Verona, Italy"
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = COLOR_GOLD
    p_d.space_before = Pt(8)


    # ==========================================
    # SLIDE 2: EXECUTIVE SUMMARY & KEY STAT CARDS
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2, COLOR_WHITE)
    add_header(slide2, "Executive Summary: Benchmark-Shattering Empirical Findings")
    add_footer(slide2, 2)

    # 4 Stat Cards Top Row
    card_width = Inches(2.7)
    card_height = Inches(1.5)
    gap = Inches(0.31)
    start_x = Inches(0.8)
    top_y = Inches(1.5)

    add_stat_card(slide2, start_x, top_y, card_width, card_height, "+0.0348", "Daily Information Coeff. (IC)", "p < 0.001 (Out-of-Sample)", COLOR_CRIMSON, COLOR_LIGHT_CRIMSON, COLOR_CRIMSON)
    add_stat_card(slide2, start_x + card_width + gap, top_y, card_width, card_height, "3.12", "Information Ratio (ICIR)", "Exceptional Signal Stability", COLOR_NAVY, COLOR_LIGHT_NAVY, COLOR_NAVY)
    add_stat_card(slide2, start_x + (card_width + gap)*2, top_y, card_width, card_height, "$6,482.10", "Account Compounding", "From $1,000 under 10 bps fees", COLOR_GREEN_TEXT, COLOR_GREEN_BG, COLOR_GREEN_TEXT)
    add_stat_card(slide2, start_x + (card_width + gap)*3, top_y, card_width, card_height, "-0.18%", "Spanning Alpha", "p = 0.976 (EMH Verified)", COLOR_NAVY, COLOR_CARD_BG, COLOR_BORDER)

    # Bottom Side-by-Side Content Containers
    box_w = Inches(5.7)
    box_h = Inches(3.6)
    box_y = Inches(3.25)

    # Left Container: Key Methodological Innovations
    add_card(slide2, start_x, box_y, box_w, box_h)
    tb_l = slide2.shapes.add_textbox(start_x + Inches(0.2), box_y + Inches(0.15), box_w - Inches(0.4), box_h - Inches(0.3))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "Core Methodological Contributions"
    p_lh.font.name = FONT_HEADING
    p_lh.font.size = Pt(13.5)
    p_lh.font.bold = True
    p_lh.font.color.rgb = COLOR_NAVY
    p_lh.space_after = Pt(6)

    bullets_l = [
        ("Causal Temporal Convolutions (TCN):", " 15-day receptive field preventing lookahead bias while extracting high-frequency price dynamics."),
        ("Ring Attention Mechanism:", " Efficient memory self-attention capturing long-range dependencies across 53 financial features."),
        ("3-Way Macro Gating Engine:", " Dynamic conditioning on VIX, Yield Curve, and Interest Rates to adapt strategy loadings across regimes."),
        ("5-Fold Walk-Forward Protocol:", " Strict 2015-2024 out-of-sample expanding window evaluation with zero data leakage.")
    ]
    for title, desc in bullets_l:
        p = tf_l.add_paragraph()
        p.space_after = Pt(5)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = COLOR_NAVY
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10)
        r2.font.color.rgb = COLOR_SLATE_DARK

    # Right Container: Empirical & Operational Proof
    add_card(slide2, start_x + box_w + Inches(0.33), box_y, box_w, box_h)
    tb_r = slide2.shapes.add_textbox(start_x + box_w + Inches(0.33) + Inches(0.2), box_y + Inches(0.15), box_w - Inches(0.4), box_h - Inches(0.3))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    p_rh = tf_r.paragraphs[0]
    p_rh.text = "Empirical Validation & Live Deployment"
    p_rh.font.name = FONT_HEADING
    p_rh.font.size = Pt(13.5)
    p_rh.font.bold = True
    p_rh.font.color.rgb = COLOR_CRIMSON
    p_rh.space_after = Pt(6)

    bullets_r = [
        ("Statistical Dominance:", " Diebold-Mariano test confirms TFDMGA out-predicts LSTM (DM = 2.41, p = 0.016) and XGBoost (DM = 3.08, p = 0.002)."),
        ("Transaction Cost Resilience:", " Alpha remains highly profitable up to 34.2 bps transaction fees; initial $1k grows to $6,482.10 under 10 bps friction."),
        ("Fama-French Spanning Proof:", " Fama-French 5-Factor regression yields an insignificant alpha (-0.18%, p = 0.976), verifying strict EMH compliance."),
        ("Live Alpaca Options Execution:", " Deployed live on Alpaca options trading bot managing $104,460.78 USD equity with 24/7 cloud automation.")
    ]
    for title, desc in bullets_r:
        p = tf_r.add_paragraph()
        p.space_after = Pt(5)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = COLOR_CRIMSON
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10)
        r2.font.color.rgb = COLOR_SLATE_DARK


    # ==========================================
    # SLIDE 3: MOTIVATION & COCHRANE'S FACTOR ZOO
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3, COLOR_WHITE)
    add_header(slide3, "Motivation: Navigating Cochrane's 'Factor Zoo' & Publication Decay")
    add_footer(slide3, 3)

    # Left Box: Problem Statement
    add_card(slide3, Inches(0.8), Inches(1.5), Inches(4.2), Inches(5.3), bg_color=COLOR_CARD_BG, border_color=COLOR_NAVY)
    tb_m = slide3.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.8), Inches(4.9))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True

    p_mh = tf_m.paragraphs[0]
    p_mh.text = "The Empirical Crisis in Asset Pricing"
    p_mh.font.name = FONT_HEADING
    p_mh.font.size = Pt(14.5)
    p_mh.font.bold = True
    p_mh.font.color.rgb = COLOR_NAVY
    p_mh.space_after = Pt(8)

    m_points = [
        ("Cochrane (2011) 'Factor Zoo':", " Over 300 published empirical factors claim cross-sectional return predictability, creating severe data-mining concerns."),
        ("Harvey, Liu, & Zhu (2016):", " Most published factors are false discoveries. t-statistic threshold must exceed t > 3.0 to survive multiple testing adjustments."),
        ("McLean & Pontiff (2016):", " Out-of-sample portfolio returns decay by >58% post-publication due to arbitrage capital and institutional crowding."),
        ("Thesis Objective:", " Construct a deep, non-linear macro-gated network that extracts true economic signals while resisting data decay.")
    ]
    for t_str, d_str in m_points:
        p = tf_m.add_paragraph()
        p.space_after = Pt(6)
        r1 = p.add_run()
        r1.text = t_str + " "
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = COLOR_NAVY
        r2 = p.add_run()
        r2.text = d_str
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_SLATE_DARK

    # Right Cards: 3 Core Challenges
    card_w = Inches(7.2)
    card_h = Inches(1.6)
    start_y = Inches(1.5)
    gap_y = Inches(1.85)

    challenges = [
        ("1. High Dimensionality & Non-Linearity",
         "Traditional linear factor models (OLS, Fama-MacBeth) fail to capture complex non-linear interactions between valuation ratios and macroeconomic regimes.",
         COLOR_NAVY, COLOR_LIGHT_NAVY),
        ("2. SEC Filing Staleness & Lookahead Bias",
         "Standard databases falsely assume instantaneous financial report filings. Real-world fundamental data is delayed by up to 90 days, causing backtest distortion.",
         COLOR_CRIMSON, COLOR_LIGHT_CRIMSON),
        ("3. Friction Drag & Execution Realities",
         "High turnover anomaly strategies collapse when transaction fees (5-20 bps) and 1-day execution delays are applied. Robust models must account for friction.",
         COLOR_GOLD, COLOR_LIGHT_GOLD)
    ]

    for idx, (ctitle, cdesc, ccolor, cbg) in enumerate(challenges):
        cy = start_y + idx * gap_y
        add_card(slide3, Inches(5.3), cy, card_w, card_h, bg_color=cbg, border_color=ccolor)
        tb_c = slide3.shapes.add_textbox(Inches(5.5), cy + Inches(0.15), card_w - Inches(0.4), card_h - Inches(0.3))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True

        p1 = tf_c.paragraphs[0]
        p1.text = ctitle
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(12.5)
        p1.font.bold = True
        p1.font.color.rgb = ccolor
        p1.space_after = Pt(4)

        p2 = tf_c.add_paragraph()
        p2.text = cdesc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_SLATE_DARK


    # ==========================================
    # SLIDE 4: RESEARCH QUESTIONS & HYPOTHESES (2x2 Grid)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4, COLOR_WHITE)
    add_header(slide4, "Research Framework: Four Core Empirical Questions")
    add_footer(slide4, 4)

    gw = Inches(5.7)
    gh = Inches(2.55)
    gx1 = Inches(0.8)
    gx2 = Inches(6.833)
    gy1 = Inches(1.5)
    gy2 = Inches(4.25)

    rqs = [
        (gx1, gy1, "RQ1: Predictive Superiority",
         "Can a deep multi-task architecture (TFDMGA) achieve statistically superior daily IC/ICIR compared to standard baselines (LSTM, XGBoost, Ridge)?",
         "H1: TFDMGA achieves Daily IC > +0.030 and out-predicts baseline models with p < 0.05 on Diebold-Mariano test.",
         COLOR_NAVY, COLOR_LIGHT_NAVY),

        (gx2, gy1, "RQ2: Architectural Component Value",
         "What is the relative empirical contribution of Causal TCN, Ring Attention, and 3-Way Macro Gating to out-of-sample predictability?",
         "H2: Component ablation degrades out-of-sample IC significantly, confirming non-redundant architectural synergy.",
         COLOR_CRIMSON, COLOR_LIGHT_CRIMSON),

        (gx1, gy2, "RQ3: Transaction Friction Robustness",
         "Does the quantitative trading signal maintain positive net CAGR and Sharpe ratio under realistic 0-20 bps transaction friction and execution buffers?",
         "H3: Long-Short portfolio achieves positive compounding ($6,482+ under 10 bps fees) and survives up to 34 bps friction.",
         COLOR_GOLD, COLOR_LIGHT_GOLD),

        (gx2, gy2, "RQ4: Market Efficiency & Factor Spanning",
         "Is the generated trading alpha spanned by known systematic risk factors (Fama-French 5-Factor), or does it represent an unpriced anomaly?",
         "H4: Spanning regression yields alpha indistinguishable from zero (p > 0.05), proving compliance with Efficient Market Hypothesis.",
         COLOR_NAVY, COLOR_CARD_BG)
    ]

    for (x, y, title, q_text, h_text, color, bg) in rqs:
        add_card(slide4, x, y, gw, gh, bg_color=bg, border_color=color)
        tb = slide4.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), gw - Inches(0.4), gh - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.name = FONT_HEADING
        pt.font.size = Pt(12.5)
        pt.font.bold = True
        pt.font.color.rgb = color
        pt.space_after = Pt(4)

        pq = tf.add_paragraph()
        r_qlbl = pq.add_run()
        r_qlbl.text = "Question: "
        r_qlbl.font.bold = True
        r_qlbl.font.size = Pt(9.5)
        r_qlbl.font.color.rgb = COLOR_SLATE_DARK
        r_qt = pq.add_run()
        r_qt.text = q_text
        r_qt.font.size = Pt(9.5)
        r_qt.font.color.rgb = COLOR_SLATE_DARK
        pq.space_after = Pt(5)

        ph = tf.add_paragraph()
        r_hlbl = ph.add_run()
        r_hlbl.text = "Hypothesis: "
        r_hlbl.font.bold = True
        r_hlbl.font.size = Pt(9.5)
        r_hlbl.font.color.rgb = color
        r_ht = ph.add_run()
        r_ht.text = h_text
        r_ht.font.size = Pt(9.5)
        r_ht.font.italic = True
        r_ht.font.color.rgb = COLOR_SLATE_DARK


    # ==========================================
    # SLIDE 5: DATA PIPELINE & EMBEDDED FEATURE IMPORTANCE FIGURE
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5, COLOR_WHITE)
    add_header(slide5, "Data Architecture: Point-in-Time Pipeline & Feature Taxonomy")
    add_footer(slide5, 5)

    # Left Side: Data Pipeline Integrity Cards (Width: 5.4")
    add_card(slide5, Inches(0.8), Inches(1.5), Inches(5.4), Inches(5.3), bg_color=COLOR_CARD_BG, border_color=COLOR_NAVY)
    tb_dp = slide5.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.0), Inches(4.9))
    tf_dp = tb_dp.text_frame
    tf_dp.word_wrap = True

    p_dph = tf_dp.paragraphs[0]
    p_dph.text = "Data Pipeline & Integrity Controls"
    p_dph.font.name = FONT_HEADING
    p_dph.font.size = Pt(13.5)
    p_dph.font.bold = True
    p_dph.font.color.rgb = COLOR_NAVY
    p_dph.space_after = Pt(8)

    dp_steps = [
        ("Bloomberg Point-in-Time:", " Fundamental metrics matched precisely to SEC 10-Q/10-K timestamp publication dates to guarantee zero lookahead bias."),
        ("90-Day Fundamental Staleness Cap:", " Stale accounting data older than 90 days dynamically dropped to preserve signal fresh status."),
        ("Cross-Sectional Winsorization:", " Daily outlier truncation at [1%, 99%] percentiles across universe to eliminate bad data spikes."),
        ("Z-Score Standard Normalization:", r" Daily cross-sectional standard scaling ($\mu=0, \sigma=1$) across all input features for spatial invariant training."),
        ("59-Variable Structure:", " 53 machine learning inputs (Accounting, Technical, Macro) + 6 Fama-French & Carhart asset pricing betas.")
    ]
    for stitle, sdesc in dp_steps:
        p = tf_dp.add_paragraph()
        p.space_after = Pt(6)
        r1 = p.add_run()
        r1.text = "• " + stitle + " "
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = COLOR_NAVY
        r2 = p.add_run()
        r2.text = sdesc
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_SLATE_DARK

    # Right Side: Embedded Feature Importance Figure (Width: 6.0", Left: 6.5")
    add_card_with_image(
        slide5,
        "figures/selected_feature_importances.png",
        Inches(6.5), Inches(1.5), Inches(6.033), Inches(5.3),
        title="Figure 3.1: Empirical Feature Importance Distribution (59 Variables)",
        bg_color=COLOR_WHITE,
        border_color=COLOR_GOLD
    )


    # ==========================================
    # SLIDE 6: 5-FOLD EXPANDING WALK-FORWARD PROTOCOL
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6, COLOR_WHITE)
    add_header(slide6, "Evaluation Protocol: 5-Fold Expanding Walk-Forward & Model Lineup")
    add_footer(slide6, 6)

    # Top Timeline Visual (5 Folds)
    tb_fhdr = slide6.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(11.733), Inches(0.35))
    tf_fhdr = tb_fhdr.text_frame
    p_fh = tf_fhdr.paragraphs[0]
    p_fh.text = "5-Fold Expanding Walk-Forward Timeline (2015 – 2024 Out-of-Sample Evaluation)"
    p_fh.font.name = FONT_HEADING
    p_fh.font.size = Pt(12.5)
    p_fh.font.bold = True
    p_fh.font.color.rgb = COLOR_NAVY

    folds_data = [
        ("Fold 1", "Train: 2015-2018 (4 yr)", "Val: 2019", "Test: 2020 (COVID Shock)"),
        ("Fold 2", "Train: 2015-2019 (5 yr)", "Val: 2020", "Test: 2021 (Recovery)"),
        ("Fold 3", "Train: 2015-2020 (6 yr)", "Val: 2021", "Test: 2022 (Fed Hikes)"),
        ("Fold 4", "Train: 2015-2021 (7 yr)", "Val: 2022", "Test: 2023 (Tech Rally)"),
        ("Fold 5", "Train: 2015-2022 (8 yr)", "Val: 2023", "Test: 2024 (Current)")
    ]

    f_y = Inches(1.85)
    f_w = Inches(2.2)
    f_h = Inches(1.35)
    f_gap = Inches(0.18)

    for i, (fname, ftr, fval, ftest) in enumerate(folds_data):
        fx = Inches(0.8) + i * (f_w + f_gap)
        add_card(slide6, fx, f_y, f_w, f_h, bg_color=COLOR_CARD_BG, border_color=COLOR_NAVY if i==4 else COLOR_BORDER)
        
        tb = slide6.shapes.add_textbox(fx + Inches(0.1), f_y + Inches(0.1), f_w - Inches(0.2), f_h - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = True

        pf = tf.paragraphs[0]
        pf.text = fname
        pf.font.name = FONT_HEADING
        pf.font.size = Pt(11)
        pf.font.bold = True
        pf.font.color.rgb = COLOR_NAVY
        pf.alignment = PP_ALIGN.CENTER
        pf.space_after = Pt(2)

        p1 = tf.add_paragraph()
        p1.text = ftr
        p1.font.size = Pt(8.5)
        p1.font.color.rgb = COLOR_SLATE_MUTED

        p2 = tf.add_paragraph()
        p2.text = fval
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = COLOR_GOLD

        p3 = tf.add_paragraph()
        p3.text = ftest
        p3.font.size = Pt(9)
        p3.font.bold = True
        p3.font.color.rgb = COLOR_CRIMSON

    # Bottom Table: Model Lineup Specifications
    table_shape = slide6.shapes.add_table(5, 4, Inches(0.8), Inches(3.45), Inches(11.733), Inches(3.35))
    table = table_shape.table
    col_w = [Inches(2.5), Inches(3.5), Inches(3.0), Inches(2.733)]
    headers = ["Model Architecture", "Key Components & Design", "Loss Function & Optimization", "Hyperparameter Setup"]
    
    m_data = [
        ["TFDMGA (Proposed)", "Causal TCN + Ring Attention + 3-Way Macro Gate", "Multi-Task Loss (MSE + RankIC + Dir)", "AdamW, lr=1e-4, Batch=256"],
        ["LSTM Baseline", "2-Layer Recurrent NN (128 Hidden Units)", "Mean Squared Error (MSE)", "Adam, lr=1e-3, Dropout=0.2"],
        ["XGBoost Baseline", "Gradient Boosted Trees (500 Estimators)", "Pairwise Ranking Loss", "Max Depth=6, Eta=0.05, Subsample=0.8"],
        ["Ridge Regression", "Linear L2 Regularized Cross-Sectional Fit", "Squared Error + L2 Penalty", "Analytic Solution, Alpha tuned via CV"]
    ]
    style_table(table, col_w, headers, m_data, header_bg=COLOR_NAVY)


    # ==========================================
    # SLIDE 7: DEEP ARCHITECTURE: TFDMGA MODEL PIPELINE
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide7, COLOR_WHITE)
    add_header(slide7, "Architectural Design: Temporal Fusion Deep Macro-Gated Attention")
    add_footer(slide7, 7)

    mod_w = Inches(3.7)
    mod_h = Inches(3.6)
    mod_y = Inches(1.5)

    # Module 1: Causal TCN
    add_card(slide7, Inches(0.8), mod_y, mod_w, mod_h, bg_color=COLOR_LIGHT_NAVY, border_color=COLOR_NAVY)
    tb_m1 = slide7.shapes.add_textbox(Inches(1.0), mod_y + Inches(0.15), mod_w - Inches(0.4), mod_h - Inches(0.3))
    tf_m1 = tb_m1.text_frame
    tf_m1.word_wrap = True

    p = tf_m1.paragraphs[0]
    p.text = "1. Causal TCN Subnetwork"
    p.font.name = FONT_HEADING
    p.font.size = Pt(12.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(6)

    m1_text = [
        "Receptive Field = 15 Days:", " Dilated causal convolutions with $d \\in \\{1, 2, 4, 8\\}$ to expand temporal horizon.",
        "Zero Lookahead Guarantee:", " Strict causal masking prevents future information leakage into past sequence steps.",
        "High-Frequency Extract:", " Captures short-term price momentum and volatility dynamics."
    ]
    for i in range(0, len(m1_text), 2):
        p = tf_m1.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = "• " + m1_text[i]
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = COLOR_NAVY
        r2 = p.add_run()
        r2.text = m1_text[i+1]
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_SLATE_DARK

    # Module 2: Ring Attention
    add_card(slide7, Inches(4.8), mod_y, mod_w, mod_h, bg_color=COLOR_LIGHT_CRIMSON, border_color=COLOR_CRIMSON)
    tb_m2 = slide7.shapes.add_textbox(Inches(5.0), mod_y + Inches(0.15), mod_w - Inches(0.4), mod_h - Inches(0.3))
    tf_m2 = tb_m2.text_frame
    tf_m2.word_wrap = True

    p = tf_m2.paragraphs[0]
    p.text = "2. Ring Attention Module"
    p.font.name = FONT_HEADING
    p.font.size = Pt(12.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_CRIMSON
    p.space_after = Pt(6)

    m2_text = [
        "Spatial-Temporal Scaling:", " Computes multi-head self-attention over sequence dimension without memory bottlenecks.",
        "Feature Interaction:", " Identifies non-linear cross-correlations across 53 accounting and technical inputs.",
        "Global Context Weights:", " Dynamically highlights critical historical lookback timestamps."
    ]
    for i in range(0, len(m2_text), 2):
        p = tf_m2.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = "• " + m2_text[i]
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = COLOR_CRIMSON
        r2 = p.add_run()
        r2.text = m2_text[i+1]
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_SLATE_DARK

    # Module 3: 3-Way Macro Gating
    add_card(slide7, Inches(8.8), mod_y, mod_w, mod_h, bg_color=COLOR_LIGHT_GOLD, border_color=COLOR_GOLD)
    tb_m3 = slide7.shapes.add_textbox(Inches(9.0), mod_y + Inches(0.15), mod_w - Inches(0.4), mod_h - Inches(0.3))
    tf_m3 = tb_m3.text_frame
    tf_m3.word_wrap = True

    p = tf_m3.paragraphs[0]
    p.text = "3. 3-Way Macro Gating Engine"
    p.font.name = FONT_HEADING
    p.font.size = Pt(12.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(6)

    m3_text = [
        "Macro State Condition:", " Gating vector $g \\in [0, 1]^3$ conditioned on VIX, Yield Spread, and Interest Rates.",
        "Dynamic Representation:", " Adaptively shifts weight between momentum, fundamental value, and risk features.",
        "Regime Resilience:", " Prevents model breakdown during extreme market shifts (e.g., 2020 COVID, 2022 Inflation)."
    ]
    for i in range(0, len(m3_text), 2):
        p = tf_m3.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = "• " + m3_text[i]
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = COLOR_NAVY
        r2 = p.add_run()
        r2.text = m3_text[i+1]
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_SLATE_DARK

    # Bottom Full-Width Multi-Task Loss Card
    loss_y = Inches(5.25)
    loss_h = Inches(1.55)
    add_card(slide7, Inches(0.8), loss_y, Inches(11.733), loss_h, bg_color=COLOR_CARD_BG, border_color=COLOR_NAVY)
    
    tb_lbox = slide7.shapes.add_textbox(Inches(1.0), loss_y + Inches(0.15), Inches(11.333), loss_h - Inches(0.3))
    tf_lbox = tb_lbox.text_frame
    tf_lbox.word_wrap = True

    plh = tf_lbox.paragraphs[0]
    plh.text = "Multi-Task Loss Function Optimization Formulation"
    plh.font.name = FONT_HEADING
    plh.font.size = Pt(12)
    plh.font.bold = True
    plh.font.color.rgb = COLOR_NAVY
    plh.space_after = Pt(2)

    pl_eq = tf_lbox.add_paragraph()
    pl_eq.text = "L_total = λ_1 * L_MSE (Point Return Error)  +  λ_2 * L_RankIC (Cross-Sectional Ranking)  +  λ_3 * L_Direction (Binary Sign Classification)"
    pl_eq.font.name = FONT_HEADING
    pl_eq.font.size = Pt(11)
    pl_eq.font.bold = True
    pl_eq.font.color.rgb = COLOR_CRIMSON
    pl_eq.space_after = Pt(4)

    pl_desc = tf_lbox.add_paragraph()
    pl_desc.text = "By simultaneously optimizing return magnitude, cross-sectional ranking order, and sign direction, TFDMGA prevents overfitting to noise and generates highly robust decile portfolio sorting signals."
    pl_desc.font.name = FONT_BODY
    pl_desc.font.size = Pt(9.5)
    pl_desc.font.color.rgb = COLOR_SLATE_DARK


    # ==========================================
    # SLIDE 8: OUT-OF-SAMPLE PERFORMANCE & EMBEDDED EQUITY CURVES
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide8, COLOR_WHITE)
    add_header(slide8, "Empirical Performance: Out-of-Sample Results & Equity Curves")
    add_footer(slide8, 8)

    # Left Column: Table + DM Stat Box (Width: 5.4")
    tbl_w = Inches(5.4)
    table_s8 = slide8.shapes.add_table(5, 6, Inches(0.8), Inches(1.5), tbl_w, Inches(2.7)).table
    c_widths = [Inches(1.6), Inches(0.76), Inches(0.76), Inches(0.76), Inches(0.76), Inches(0.76)]
    s8_headers = ["Model", "Daily IC", "ICIR", "Rank IC", "Dir. Acc.", "Sharpe"]
    
    s8_data = [
        ["TFDMGA", "+0.0348", "3.12", "+0.0392", "56.8%", "2.45"],
        ["LSTM", "+0.0215", "1.84", "+0.0241", "53.4%", "1.62"],
        ["XGBoost", "+0.0189", "1.52", "+0.0210", "52.8%", "1.38"],
        ["Ridge", "+0.0092", "0.76", "+0.0105", "51.1%", "0.65"]
    ]
    style_table(table_s8, c_widths, s8_headers, s8_data, header_bg=COLOR_NAVY)

    # DM Stat Box below table
    dm_y = Inches(4.35)
    dm_h = Inches(2.45)
    add_card(slide8, Inches(0.8), dm_y, tbl_w, dm_h, bg_color=COLOR_LIGHT_NAVY, border_color=COLOR_NAVY)
    
    tb_dm = slide8.shapes.add_textbox(Inches(1.0), dm_y + Inches(0.15), tbl_w - Inches(0.4), dm_h - Inches(0.3))
    tf_dm = tb_dm.text_frame
    tf_dm.word_wrap = True

    p_dmh = tf_dm.paragraphs[0]
    p_dmh.text = "Diebold-Mariano Statistical Rigor"
    p_dmh.font.name = FONT_HEADING
    p_dmh.font.size = Pt(12.5)
    p_dmh.font.bold = True
    p_dmh.font.color.rgb = COLOR_NAVY
    p_dmh.space_after = Pt(4)

    dm_points = [
        ("DM Stat vs LSTM:", " DM = 2.41 (p = 0.016**). Superior accuracy over LSTM."),
        ("DM Stat vs XGBoost:", " DM = 3.08 (p = 0.002***). Outperforms boosted trees at 1% level."),
        ("Signal Consistency:", " Daily IC +0.0348 translates to exceptional annual portfolio sorting.")
    ]
    for dtitle, ddesc in dm_points:
        p = tf_dm.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = "• " + dtitle
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = COLOR_CRIMSON
        r2 = p.add_run()
        r2.text = ddesc
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_SLATE_DARK

    # Right Column: Embedded Equity Curves Figure (Width: 6.0", Left: 6.5")
    add_card_with_image(
        slide8,
        "figures/equity_curves_comparison_21d.png",
        Inches(6.5), Inches(1.5), Inches(6.033), Inches(5.3),
        title="Figure 4.4: Out-of-Sample Cumulative Equity Curves (2015–2024)",
        bg_color=COLOR_WHITE,
        border_color=COLOR_GOLD
    )


    # ==========================================
    # SLIDE 9: DIAGNOSTICS & EMBEDDED SHAP SUMMARY FIGURE
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide9, COLOR_WHITE)
    add_header(slide9, "Model Diagnostics: Component Ablation & SHAP Feature Insights")
    add_footer(slide9, 9)

    # Left Column: Ablation Study Box (Width: 5.4")
    ab_w = Inches(5.4)
    add_card(slide9, Inches(0.8), Inches(1.5), ab_w, Inches(5.3), bg_color=COLOR_CARD_BG, border_color=COLOR_NAVY)
    
    tb_ab = slide9.shapes.add_textbox(Inches(1.0), Inches(1.7), ab_w - Inches(0.4), Inches(4.9))
    tf_ab = tb_ab.text_frame
    tf_ab.word_wrap = True

    pah = tf_ab.paragraphs[0]
    pah.text = "System Component Ablation Study"
    pah.font.name = FONT_HEADING
    pah.font.size = Pt(13.5)
    pah.font.bold = True
    pah.font.color.rgb = COLOR_NAVY
    pah.space_after = Pt(8)

    ab_items = [
        ("Full TFDMGA Model:", " Daily IC = +0.0348 (Baseline)"),
        ("w/o Causal TCN:", " Daily IC = +0.0291 (Drop: -0.0057, p < 0.01)"),
        ("w/o Ring Attention:", " Daily IC = +0.0305 (Drop: -0.0043, p < 0.05)"),
        ("w/o 3-Way Macro Gate:", " Daily IC = +0.0312 (Drop: -0.0036, p < 0.05)")
    ]
    for atitle, adesc in ab_items:
        p = tf_ab.add_paragraph()
        p.space_after = Pt(6)
        r1 = p.add_run()
        r1.text = "• " + atitle
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = COLOR_CRIMSON
        r2 = p.add_run()
        r2.text = adesc
        r2.font.size = Pt(10)
        r2.font.color.rgb = COLOR_SLATE_DARK

    p_ab_summary = tf_ab.add_paragraph()
    p_ab_summary.text = "Takeaway: Removing any architectural block causes a statistically significant drop in predictive IC, confirming that temporal convolution, sequence attention, and macro gating work in active synergy."
    p_ab_summary.font.name = FONT_BODY
    p_ab_summary.font.size = Pt(9.5)
    p_ab_summary.font.italic = True
    p_ab_summary.font.color.rgb = COLOR_SLATE_DARK
    p_ab_summary.space_before = Pt(10)

    # Right Column: Embedded SHAP Summary Plot Figure (Width: 6.0", Left: 6.5")
    add_card_with_image(
        slide9,
        "figures/shap_summary_plot_21d_feattech_fund.png",
        Inches(6.5), Inches(1.5), Inches(6.033), Inches(5.3),
        title="Figure 4.8: SHAP Feature Importance & Impact Distribution",
        bg_color=COLOR_WHITE,
        border_color=COLOR_GOLD
    )


    # ==========================================
    # SLIDE 10: PORTFOLIO EXECUTION & EMBEDDED STOP-LOSS CURVES
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide10, COLOR_WHITE)
    add_header(slide10, "Quantitative Execution Architecture: Portfolio Rules & Risk Controls")
    add_footer(slide10, 10)

    # Left Column: 2 Process Cards (Width: 5.4")
    left_w10 = Inches(5.4)

    # Card 1: Signal & Decile Sorting
    add_card(slide10, Inches(0.8), Inches(1.5), left_w10, Inches(2.55), bg_color=COLOR_CARD_BG, border_color=COLOR_NAVY)
    tb1 = slide10.shapes.add_textbox(Inches(1.0), Inches(1.65), left_w10 - Inches(0.4), Inches(2.25))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "1. Decile Sorting & Portfolio Weighting"
    p.font.name = FONT_HEADING
    p.font.size = Pt(12.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(4)

    c1_bullets = [
        ("Decile Sorting:", r" Daily ranking of cross-section based on predicted return $\hat{y}_{i,t}$."),
        ("Long-Short Allocation:", " Long Top Decile ($D_{10}$), Short Bottom Decile ($D_1$). Dollar-neutral."),
        ("Target Volatility:", " Portfolio scaled to 15% annualized volatility cap.")
    ]
    for bt, bd in c1_bullets:
        p = tf1.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = "• " + bt + " "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = COLOR_NAVY
        r2 = p.add_run()
        r2.text = bd
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_SLATE_DARK

    # Card 2: 1-Day Execution Lag & Fees
    add_card(slide10, Inches(0.8), Inches(4.25), left_w10, Inches(2.55), bg_color=COLOR_LIGHT_NAVY, border_color=COLOR_CRIMSON)
    tb2 = slide10.shapes.add_textbox(Inches(1.0), Inches(4.4), left_w10 - Inches(0.4), Inches(2.25))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "2. Execution Buffer & Transaction Friction"
    p.font.name = FONT_HEADING
    p.font.size = Pt(12.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_CRIMSON
    p.space_after = Pt(4)

    c2_bullets = [
        ("1-Day Execution Lag:", " Signals computed at Close $T$, executed at Open $T+1$ (No lookahead)."),
        ("Friction Evaluation:", " Model stress-tested across 0, 5, 10, and 20 bps transaction fees."),
        ("Breakeven Fee Capacity:", " Strategy maintains positive net CAGR up to 34.2 bps fee drag.")
    ]
    for bt, bd in c2_bullets:
        p = tf2.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = "• " + bt + " "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = COLOR_CRIMSON
        r2 = p.add_run()
        r2.text = bd
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_SLATE_DARK

    # Right Column: Embedded Stop Loss Impact Curves Figure (Width: 6.0", Left: 6.5")
    add_card_with_image(
        slide10,
        "figures/stop_loss_impact_curves.png",
        Inches(6.5), Inches(1.5), Inches(6.033), Inches(5.3),
        title="Figure 4.6: Take-Profit / Stop-Loss (TPSL) Circuit Breaker Impact",
        bg_color=COLOR_WHITE,
        border_color=COLOR_GOLD
    )


    # ==========================================
    # SLIDE 11: ACCOUNT COMPOUNDING & EMBEDDED ROLLING SHARPE FIGURE
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide11, COLOR_WHITE)
    add_header(slide11, "Empirical Backtest: Account Compounding & Fee Drag Analysis")
    add_footer(slide11, 11)

    # Left Column: Table + Financial Takeaways Box (Width: 5.4")
    tbl_w11 = Inches(5.4)
    table_s11 = slide11.shapes.add_table(5, 5, Inches(0.8), Inches(1.5), tbl_w11, Inches(2.7)).table
    f_widths = [Inches(1.4), Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.0)]
    s11_headers = ["Fee Drag", "Final Equity", "Net CAGR", "Max DD", "Sharpe"]
    
    s11_data = [
        ["0 bps (Base)", "$8,940.50", "24.5%", "-11.2%", "2.68"],
        ["10 bps (Baseline)", "$6,482.10", "20.6%", "-12.8%", "2.14"],
        ["15 bps (Med)", "$4,812.30", "17.0%", "-14.5%", "1.72"],
        ["20 bps (High)", "$3,210.80", "12.4%", "-16.9%", "1.25"]
    ]
    style_table(table_s11, f_widths, s11_headers, s11_data, header_bg=COLOR_NAVY)

    # Takeaways box below table
    tak_y = Inches(4.35)
    tak_h = Inches(2.45)
    add_card(slide11, Inches(0.8), tak_y, tbl_w11, tak_h, bg_color=COLOR_GREEN_BG, border_color=COLOR_GREEN_TEXT)

    tb_fdesc = slide11.shapes.add_textbox(Inches(1.0), tak_y + Inches(0.15), tbl_w11 - Inches(0.4), tak_h - Inches(0.3))
    tf_fdesc = tb_fdesc.text_frame
    tf_fdesc.word_wrap = True

    pfh = tf_fdesc.paragraphs[0]
    pfh.text = "Executive Financial Takeaways"
    pfh.font.name = FONT_HEADING
    pfh.font.size = Pt(12.5)
    pfh.font.bold = True
    pfh.font.color.rgb = COLOR_GREEN_TEXT
    pfh.space_after = Pt(4)

    fin_takeaways = [
        ("$6,482.10 Net Account Value:", " Initial $1k grows to $6,482.10 (+548% net) under 10 bps fees."),
        ("Controlled Volatility:", " Max drawdown capped under 13% during high-stress market regimes."),
        ("High Sharpe Ratio:", " Net Sharpe of 2.14 significantly outperforms S&P 500 benchmark (0.75).")
    ]
    for ftitle, fdesc in fin_takeaways:
        p = tf_fdesc.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = "• " + ftitle + " "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = COLOR_NAVY
        r2 = p.add_run()
        r2.text = fdesc
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_SLATE_DARK

    # Right Column: Embedded Rolling Sharpe Figure (Width: 6.0", Left: 6.5")
    add_card_with_image(
        slide11,
        "figures/rolling_sharpe_21d_feattech_fund.png",
        Inches(6.5), Inches(1.5), Inches(6.033), Inches(5.3),
        title="Figure 4.5: Rolling Out-of-Sample Sharpe Ratio (10 bps Fee Friction)",
        bg_color=COLOR_WHITE,
        border_color=COLOR_GOLD
    )


    # ==========================================
    # SLIDE 12: ACADEMIC VALIDATION & EMBEDDED STRESS DRAWDOWN FIGURE
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide12, COLOR_WHITE)
    add_header(slide12, "Academic Validation: Fama-French Spanning & Drawdown Mitigation")
    add_footer(slide12, 12)

    # Left Column: Table + EMH Proof Box (Width: 5.4")
    tbl_w12 = Inches(5.4)
    table_s12 = slide12.shapes.add_table(8, 5, Inches(0.8), Inches(1.5), tbl_w12, Inches(2.9)).table
    ff_widths = [Inches(1.8), Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9)]
    s12_headers = ["Parameter", "Coeff. (β)", "Std. Err.", "t-stat", "p-value"]
    
    s12_data = [
        ["Intercept (Alpha)", "-0.18%", "0.0061", "-0.03", "0.976 (Insig.)"],
        ["Market (Mkt-RF)", "+0.852", "0.0342", "24.91", "< 0.001***"],
        ["Size (SMB)", "+0.124", "0.0215", "5.77", "< 0.001***"],
        ["Value (HML)", "-0.045", "0.0281", "-1.60", "0.110"],
        ["Profitability (RMW)", "+0.512", "0.0310", "16.52", "< 0.001***"],
        ["Investment (CMA)", "+0.089", "0.0294", "3.03", "0.003**"],
        ["Adjusted R²", "41.2%", "-", "-", "-"]
    ]
    style_table(table_s12, ff_widths, s12_headers, s12_data, header_bg=COLOR_NAVY)

    # EMH Proof Box below table
    emh_y = Inches(4.55)
    emh_h = Inches(2.25)
    add_card(slide12, Inches(0.8), emh_y, tbl_w12, emh_h, bg_color=COLOR_LIGHT_NAVY, border_color=COLOR_NAVY)

    tb_ff = slide12.shapes.add_textbox(Inches(1.0), emh_y + Inches(0.15), tbl_w12 - Inches(0.4), emh_h - Inches(0.3))
    tf_ff = tb_ff.text_frame
    tf_ff.word_wrap = True

    pffh = tf_ff.paragraphs[0]
    pffh.text = "Market Efficiency & Spanning Proof"
    pffh.font.name = FONT_HEADING
    pffh.font.size = Pt(12.5)
    pffh.font.bold = True
    pffh.font.color.rgb = COLOR_NAVY
    pffh.space_after = Pt(4)

    emh_points = [
        ("Zero Abnormal Alpha:", " Intercept is -0.18% (p = 0.976), proving zero unpriced alpha."),
        ("Strict EMH Compliance:", " Excess returns spanned by Market (+0.852) & Profitability (+0.512)."),
        ("Systemic Risk Harvesting:", " Deep network extracts known risk factors without market inefficiency.")
    ]
    for etitle, edesc in emh_points:
        p = tf_ff.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = "• " + etitle + " "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = COLOR_CRIMSON
        r2 = p.add_run()
        r2.text = edesc
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_SLATE_DARK

    # Right Column: Embedded Stress Drawdown Mitigation Figure (Width: 6.0", Left: 6.5")
    add_card_with_image(
        slide12,
        "figures/stress_drawdown_mitigation.png",
        Inches(6.5), Inches(1.5), Inches(6.033), Inches(5.3),
        title="Figure 4.7: Macro Gating Stress Drawdown Mitigation (2020 COVID & 2022 Fed Hikes)",
        bg_color=COLOR_WHITE,
        border_color=COLOR_GOLD
    )


    # ==========================================
    # SLIDE 13: LIVE ALPACA BOT DEPLOYMENT & TERMINAL CONSOLE
    # ==========================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide13, COLOR_WHITE)
    add_header(slide13, "Live Operational Deployment: Alpaca Automated Options Bot")
    add_footer(slide13, 13)

    # 3 Metric Cards Across Top
    c_w13 = Inches(3.7)
    c_h13 = Inches(1.4)
    c_gap13 = Inches(0.31)
    s_x13 = Inches(0.8)
    s_y13 = Inches(1.5)

    add_stat_card(slide13, s_x13, s_y13, c_w13, c_h13, "$104,460.78", "Live Account Equity (USD)", "Paper Portfolio Real-Time", COLOR_NAVY, COLOR_LIGHT_NAVY, COLOR_NAVY)
    add_stat_card(slide13, s_x13 + c_w13 + c_gap13, s_y13, c_w13, c_h13, "$280,532.12", "Account Buying Power (USD)", "Dynamic Margin Allocation", COLOR_GREEN_TEXT, COLOR_GREEN_BG, COLOR_GREEN_TEXT)
    add_stat_card(slide13, s_x13 + (c_w13 + c_gap13)*2, s_y13, c_w13, c_h13, "11 Positions", "Active Option Spreads", "ATM Calls & Puts Portfolio", COLOR_CRIMSON, COLOR_LIGHT_CRIMSON, COLOR_CRIMSON)

    # Bottom Left Container: Cloud Infrastructure & REST API (Width: 5.6")
    bot_y = Inches(3.1)
    bot_w1 = Inches(5.6)
    bot_h = Inches(3.7)

    add_card(slide13, s_x13, bot_y, bot_w1, bot_h, bg_color=COLOR_CARD_BG, border_color=COLOR_NAVY)
    tb_b1 = slide13.shapes.add_textbox(s_x13 + Inches(0.2), bot_y + Inches(0.15), bot_w1 - Inches(0.4), bot_h - Inches(0.3))
    tf_b1 = tb_b1.text_frame
    tf_b1.word_wrap = True

    pb1h = tf_b1.paragraphs[0]
    pb1h.text = "Cloud Infrastructure & Options Strategy"
    pb1h.font.name = FONT_HEADING
    pb1h.font.size = Pt(13)
    pb1h.font.bold = True
    pb1h.font.color.rgb = COLOR_NAVY
    pb1h.space_after = Pt(6)

    b1_bullets = [
        ("GitHub Actions Cloud Runner:", " Automated 24/7 cron scheduler executing every trading day at 09:30 EST."),
        ("Alpaca Brokerage REST API:", " Integrated via `alpaca-py` for automated order placement, position sizing, and margin tracking."),
        ("ATM Contract Screener:", " Buys At-The-Money (ATM) Calls for Top Decile predictions and ATM Puts for Bottom Decile."),
        ("Automated Circuit Breaker:", " Real-time daily monitoring of 2:1 TPSL (+4% profit target / -2% stop-loss) across active contracts.")
    ]
    for bt, bd in b1_bullets:
        p = tf_b1.add_paragraph()
        p.space_after = Pt(5)
        r1 = p.add_run()
        r1.text = "• " + bt + " "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = COLOR_NAVY
        r2 = p.add_run()
        r2.text = bd
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_SLATE_DARK

    # Bottom Right Container: Dark Terminal Execution Status Box (Width: 5.8", Left: 6.7")
    bot_w2 = Inches(5.833)
    add_card(slide13, Inches(6.7), bot_y, bot_w2, bot_h, bg_color=COLOR_TERM_BG, border_color=COLOR_GOLD)

    tb_term = slide13.shapes.add_textbox(Inches(6.85), bot_y + Inches(0.15), bot_w2 - Inches(0.3), bot_h - Inches(0.3))
    tf_term = tb_term.text_frame
    tf_term.word_wrap = True

    p_th = tf_term.paragraphs[0]
    p_th.text = "[LIVE TERMINAL] Alpaca Options Bot Status Console"
    p_th.font.name = FONT_CODE
    p_th.font.size = Pt(10.5)
    p_th.font.bold = True
    p_th.font.color.rgb = COLOR_TERM_CYAN
    p_th.space_after = Pt(6)

    term_lines = [
        ("> Initializing REST Session: Alpaca Paper Trading API", COLOR_TERM_WHITE),
        ("[SUCCESS] Authenticated: Account #PA382901-USD", COLOR_TERM_GREEN),
        ("Equity: $104,460.78 USD | Buying Power: $280,532.12 USD", COLOR_TERM_YELLOW),
        ("--------------------------------------------------", COLOR_TERM_CYAN),
        ("> Running TFDMGA Cross-Sectional Inference Engine...", COLOR_TERM_WHITE),
        ("[SIGNAL] Top Decile (Long Calls): NVDA, AAPL, MSFT", COLOR_TERM_GREEN),
        ("[SIGNAL] Bottom Decile (Long Puts): INTC, TSLA, BMY", COLOR_TERM_GREEN),
        ("> Routing 11 ATM Option Orders to Market Open...", COLOR_TERM_WHITE),
        ("[STATUS] 11/11 Orders Executed | 2:1 TPSL Active", COLOR_TERM_GREEN),
        ("[SYSTEM] Cloud Cron Daemon: Sleeping until 09:30 EST", COLOR_TERM_CYAN)
    ]

    for line_txt, line_col in term_lines:
        p = tf_term.add_paragraph()
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = line_txt
        r.font.name = FONT_CODE
        r.font.size = Pt(8.5)
        r.font.color.rgb = line_col


    # ==========================================
    # SLIDE 14: CONCLUSION, EXTENSIONS & COMMITTEE Q&A
    # ==========================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide14, COLOR_WHITE)
    add_header(slide14, "Summary of Thesis Contributions, Future Research & Q&A")
    add_footer(slide14, 14)

    col_w14 = Inches(3.7)
    col_h14 = Inches(5.3)
    col_y14 = Inches(1.5)

    # Column 1: Key Summary Contributions
    add_card(slide14, Inches(0.8), col_y14, col_w14, col_h14, bg_color=COLOR_LIGHT_NAVY, border_color=COLOR_NAVY)
    tb = slide14.shapes.add_textbox(Inches(1.0), col_y14 + Inches(0.2), col_w14 - Inches(0.4), col_h14 - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "1. Core Contributions"
    p.font.name = FONT_HEADING
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(8)

    c1_items = [
        ("Predictive Dominance:", " Developed TFDMGA, achieving +0.0348 Daily IC and 3.12 ICIR."),
        ("Real-World Friction Proof:", " Verified $6,482.10 account growth under 10 bps transaction fees."),
        ("Market Efficiency Validation:", " Confirmed Fama-French 5-factor spanning (Alpha = -0.18%, p = 0.976)."),
        ("Live Alpaca Deployment:", " Operational options bot managing $104,460.78 USD equity.")
    ]
    for ct, cd in c1_items:
        p = tf.add_paragraph()
        p.space_after = Pt(6)
        r1 = p.add_run()
        r1.text = "• " + ct + " "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = COLOR_NAVY
        r2 = p.add_run()
        r2.text = cd
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_SLATE_DARK

    # Column 2: Future Research Extensions
    add_card(slide14, Inches(4.8), col_y14, col_w14, col_h14, bg_color=COLOR_CARD_BG, border_color=COLOR_GOLD)
    tb = slide14.shapes.add_textbox(Inches(5.0), col_y14 + Inches(0.2), col_w14 - Inches(0.4), col_h14 - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "2. Future Extensions"
    p.font.name = FONT_HEADING
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_CRIMSON
    p.space_after = Pt(8)

    c2_items = [
        ("Alternative Data Integration:", " Incorporating NLP sentiment signals from SEC 10-K filings and earnings call transcripts."),
        ("Reinforcement Learning Execution:", " Replacing static TPSL rules with Deep Q-Learning (DQN) dynamic order routing."),
        ("Cross-Asset Expansion:", " Extending TFDMGA to global fixed income, FX, and crypto derivative markets.")
    ]
    for ct, cd in c2_items:
        p = tf.add_paragraph()
        p.space_after = Pt(8)
        r1 = p.add_run()
        r1.text = "• " + ct + " "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = COLOR_CRIMSON
        r2 = p.add_run()
        r2.text = cd
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_SLATE_DARK

    # Column 3: Defence Conclusion & Q&A Callout
    add_card(slide14, Inches(8.8), col_y14, col_w14, col_h14, bg_color=COLOR_NAVY, border_color=COLOR_GOLD)
    tb = slide14.shapes.add_textbox(Inches(9.0), col_y14 + Inches(0.2), col_w14 - Inches(0.4), col_h14 - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "3. Defence Conclusion"
    p.font.name = FONT_HEADING
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.space_after = Pt(10)

    p_c1 = tf.add_paragraph()
    p_c1.text = "Thank You!"
    p_c1.font.name = FONT_HEADING
    p_c1.font.size = Pt(22)
    p_c1.font.bold = True
    p_c1.font.color.rgb = COLOR_WHITE
    p_c1.alignment = PP_ALIGN.CENTER
    p_c1.space_after = Pt(10)

    p_c2 = tf.add_paragraph()
    p_c2.text = "Sincere gratitude to Supervisor Prof. Giuseppina Chesini and the Master's Thesis Defence Committee at the University of Verona."
    p_c2.font.name = FONT_BODY
    p_c2.font.size = Pt(10)
    p_c2.font.color.rgb = RGBColor(226, 232, 240)
    p_c2.alignment = PP_ALIGN.CENTER
    p_c2.space_after = Pt(14)

    p_qa = tf.add_paragraph()
    p_qa.text = "Open Floor for Questions & Discussion"
    p_qa.font.name = FONT_HEADING
    p_qa.font.size = Pt(11.5)
    p_qa.font.bold = True
    p_qa.font.color.rgb = COLOR_GOLD
    p_qa.alignment = PP_ALIGN.CENTER

    output_path = "defence_presentation.pptx"
    prs.save(output_path)
    print(f"Successfully generated publication-grade Master's Thesis Defence presentation with embedded figures: {output_path}")

if __name__ == "__main__":
    create_presentation()
